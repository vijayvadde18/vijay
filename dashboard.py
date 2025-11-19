# enterprise_dashboard_v3.py
# Enterprise AI Analytics — single-file mega app (v3)
# Includes: Power BI style UI + ALL Level-2 features (AI Formula, Live Connectors,
# Drilldown, Streaming, Designer, AI Cleaner, Semantic Search, Auto PPT/PDF, DQ Scoring,
# TS Anomaly Detection, Feature Eng Studio, DAX-like measures, Notebook Mode,
# Auth & RBAC, Autosave/Versioning) + previous advanced features.
#
# SECURITY: This app may execute code snippets or AI-suggested code only after explicit
# user opt-in (two toggles). Do NOT enable these on public servers without sandboxing.
#
# BEFORE RUNNING:
# pip install -r requirements.txt
# Set GROQ_API_KEY environment variable or .streamlit/secrets.toml entry.
#
# Minimal recommended requirements (some features optional):
# streamlit pandas numpy plotly scikit-learn groq ydata-profiling prophet reportlab xlsxwriter
# sentence-transformers faiss-cpu streamlit-ace ruptures python-pptx joblib sqlalchemy
#
# Run:
# streamlit run enterprise_dashboard_v3.py

import os
import io
import sys
import json
import time
import sqlite3
import base64
import joblib
import traceback
import threading
from datetime import datetime
from urllib.parse import urlparse

import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go

# ML & utilities
from sklearn.model_selection import train_test_split
from sklearn.ensemble import RandomForestRegressor, RandomForestClassifier, IsolationForest
from sklearn.metrics import mean_squared_error, r2_score, accuracy_score, classification_report
from sklearn.preprocessing import StandardScaler

# Optional libs - wrap in try/except
try:
    from prophet import Prophet
    PROPHET_AVAILABLE = True
except Exception:
    PROPHET_AVAILABLE = False

try:
    from ydata_profiling import ProfileReport
    YDATA_AVAILABLE = True
except Exception:
    YDATA_AVAILABLE = False

try:
    from groq import Groq
    GROQ_AVAILABLE = True
except Exception:
    GROQ_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    REPORTLAB_AVAILABLE = True
except Exception:
    REPORTLAB_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_AVAILABLE = True
except Exception:
    PPTX_AVAILABLE = False

# Semantic search (optional)
try:
    from sentence_transformers import SentenceTransformer
    import faiss
    STRANS_AVAILABLE = True
except Exception:
    STRANS_AVAILABLE = False

# Anomaly detection
try:
    import ruptures as rpt
    RUPTURES_AVAILABLE = True
except Exception:
    RUPTURES_AVAILABLE = False

# Notebook editor
try:
    from streamlit_ace import st_ace
    ACE_AVAILABLE = True
except Exception:
    ACE_AVAILABLE = False

# Streamlit-plotly-events for drilldown (optional)
try:
    from streamlit_plotly_events import plotly_events
    PLOTLY_EVENTS_AVAILABLE = True
except Exception:
    PLOTLY_EVENTS_AVAILABLE = False

# ---------------------------
# App config & styling (Power BI-ish)
# ---------------------------
PRIMARY = "#F0C419"   # Power BI yellow
BG = "#07111a"
CARD = "#0d2b36"
TEXT = "#E6EEF2"

st.set_page_config(page_title="Enterprise AI Analytics v3", layout="wide", initial_sidebar_state="expanded")
st.markdown(f"""
    <style>
        .reportview-container {{ background-color: {BG}; color: {TEXT}; }}
        .sidebar .sidebar-content {{ background-color: #06202a; }}
        .stButton>button {{ background-color: {PRIMARY}; color: black; border: none; }}
        .stDownloadButton>button {{ background-color: {PRIMARY}; color: black; border: none; }}
        .card {{ background: {CARD}; padding: 14px; border-radius:10px; box-shadow: 0 3px 8px rgba(0,0,0,0.45); }}
        .metric-label {{ color: #b8c2c7; font-size: 13px; }}
        .small {{ font-size:12px; color:#b8c2c7; }}
    </style>
""", unsafe_allow_html=True)

st.title("📊 Enterprise AI Analytics — All Features (v3)")

# ---------------------------
# Load Groq key & init client
# ---------------------------
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
if not GROQ_API_KEY:
    try:
        GROQ_API_KEY = st.secrets["GROQ_API_KEY"]
    except Exception:
        GROQ_API_KEY = None

groq_client = None
if GROQ_API_KEY and GROQ_AVAILABLE:
    try:
        groq_client = Groq(api_key=GROQ_API_KEY)
    except Exception:
        groq_client = None

def groq_chat(prompt: str, model: str = "llama-3.1-70b-versatile", max_tokens: int = 800):
    """Call Groq chat (graceful)."""
    if groq_client is None:
        return "Groq client not available. Set GROQ_API_KEY and install 'groq'."
    try:
        resp = groq_client.chat.completions.create(
            model=model,
            messages=[{"role":"user","content":prompt}],
            max_tokens=max_tokens
        )
        try:
            return resp.choices[0].message["content"]
        except Exception:
            return str(resp)
    except Exception as e:
        return f"Groq Error: {e}"

# ---------------------------
# Sidebar navigation + global toggles
# ---------------------------
st.sidebar.image("https://upload.wikimedia.org/wikipedia/commons/4/4a/Microsoft_Power_BI_logo.svg", width=120)
st.sidebar.markdown("## Navigation")
page = st.sidebar.radio("", [
    "Home", "Data Explorer", "Cleaning & Prep", "Auto-EDA", "Visualizations", "Drill-Down",
    "Outliers", "Forecasting", "AutoML", "AI Agent", "AI Formula Builder", "AI Data Cleaner",
    "Semantic Search", "Notebook Mode", "Feature Studio", "DAX Measures", "TS Anomaly",
    "Live Connectors", "Designer", "SQL Engine", "Export Center", "Data Quality", "Auth & RBAC",
    "Autosave & Versioning", "Streaming (Simulated)", "Deployment"
])

st.sidebar.markdown("---")
st.sidebar.markdown("### Global Settings")
ENABLE_EDA = st.sidebar.checkbox("Enable Auto-EDA (ydata_profiling)", value=False)
ENABLE_AI = st.sidebar.checkbox("Enable AI (Groq)", value=bool(groq_client))
SAFE_EXEC = st.sidebar.checkbox("Allow AI-suggested code execution (dangerous)", value=False)
CONFIRM_EXEC = st.sidebar.checkbox("Confirm execution", value=False)
ENABLE_SEMSEARCH = st.sidebar.checkbox("Enable Semantic Search (sentence-transformers+faiss)", value=False)
if ENABLE_SEMSEARCH and not STRANS_AVAILABLE:
    st.sidebar.warning("Install sentence-transformers and faiss-cpu for semantic search to work.")

# Session state initialization
if "df" not in st.session_state: st.session_state.df = None
if "cleaned_df" not in st.session_state: st.session_state.cleaned_df = None
if "models" not in st.session_state: st.session_state.models = {}
if "snapshots" not in st.session_state: st.session_state.snapshots = []
if "emb_index" not in st.session_state: st.session_state.emb_index = None
if "emb_texts" not in st.session_state: st.session_state.emb_texts = None
if "faiss_dim" not in st.session_state: st.session_state.faiss_dim = None
if "users" not in st.session_state:
    # simple sample users; in production use secure auth
    st.session_state.users = {"admin":"adminpass","analyst":"analystpass","viewer":"viewerpass"}
if "user" not in st.session_state: st.session_state.user = None

# ---------------------------
# Simple Auth (RBAC) - very basic
# ---------------------------
st.sidebar.markdown("### Authentication")
if st.session_state.user is None:
    uname = st.sidebar.text_input("Username")
    pwd = st.sidebar.text_input("Password", type="password")
    if st.sidebar.button("Login"):
        users = st.session_state.users
        if uname in users and users[uname] == pwd:
            st.session_state.user = uname
            st.sidebar.success(f"Logged in as {uname}")
        else:
            st.sidebar.error("Invalid credentials")
else:
    st.sidebar.markdown(f"**Signed in:** {st.session_state.user}")
    if st.sidebar.button("Logout"):
        st.session_state.user = None
        st.sidebar.success("Logged out")

# Helper functions
def df_to_excel_bytes(df: pd.DataFrame) -> bytes:
    towrite = io.BytesIO()
    with pd.ExcelWriter(towrite, engine="xlsxwriter") as writer:
        df.to_excel(writer, index=False)
    towrite.seek(0)
    return towrite.read()

def download_bytes(data: bytes, filename: str, mime: str):
    st.download_button(f"Download {filename}", data=data, file_name=filename, mime=mime)

def save_snapshot(df):
    os.makedirs("snapshots", exist_ok=True)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    path = f"snapshots/data_{ts}.parquet"
    df.to_parquet(path)
    st.session_state.snapshots.append(path)
    return path

# ---------------------------
# Home
# ---------------------------
if page == "Home":
    st.markdown("## Welcome to Enterprise AI Analytics v3")
    st.markdown("""
    This app includes a comprehensive set of enterprise features:
    - AI formula builder, AI data cleaner, semantic search, notebook mode, DAX-like measures
    - AutoML, forecasting, anomaly detection, PPT/PDF export, designer, streaming demo
    - Role-based auth (simple), autosave/versioning
    """)
    st.markdown("### Quick status")
    c1, c2, c3 = st.columns(3)
    with c1:
        st.markdown("<div class='card'><b>Data</b><br>", unsafe_allow_html=True)
        if st.session_state.df is None:
            st.markdown("<span class='metric-label'>No dataset loaded</span></div>", unsafe_allow_html=True)
        else:
            df = st.session_state.df
            st.markdown(f"<span class='metric-label'>{df.shape[0]:,} rows • {df.shape[1]:,} cols</span></div>", unsafe_allow_html=True)
    with c2:
        st.markdown("<div class='card'><b>AI</b><br>", unsafe_allow_html=True)
        st.markdown(f"<span class='metric-label'>{'Groq configured' if groq_client else 'Groq not configured'}</span></div>", unsafe_allow_html=True)
    with c3:
        st.markdown("<div class='card'><b>Snapshots</b><br>", unsafe_allow_html=True)
        st.markdown(f"<span class='metric-label'>{len(st.session_state.snapshots)} saved</span></div>", unsafe_allow_html=True)

# ---------------------------
# Data Explorer
# ---------------------------
elif page == "Data Explorer":
    st.header("📁 Data Explorer")
    st.markdown("Upload a CSV/XLSX, enter a URL, or load a sample dataset.")
    src = st.radio("Source", ["Sample","Upload","URL"], horizontal=True)
    if src == "Sample":
        dates = pd.date_range(end=pd.Timestamp.today(), periods=180)
        sample = pd.DataFrame({
            "date": dates,
            "region": np.random.choice(["North","South","East","West"], len(dates)),
            "product": np.random.choice(["A","B","C"], len(dates)),
            "units_sold": np.random.poisson(20, len(dates)),
            "price": np.random.uniform(10,100,len(dates))
        })
        sample["revenue"] = sample["units_sold"] * sample["price"]
        st.session_state.df = sample
        st.success("Sample dataset loaded.")
    elif src == "Upload":
        uploaded = st.file_uploader("Upload CSV/XLSX", type=["csv","xlsx"])
        if uploaded:
            try:
                if uploaded.name.lower().endswith(".csv"):
                    df = pd.read_csv(uploaded)
                else:
                    df = pd.read_excel(uploaded)
                st.session_state.df = df
                st.success("File loaded.")
            except Exception as e:
                st.error(f"Failed to read file: {e}")
    else:
        url = st.text_input("Enter data URL (CSV/JSON/XLSX)")
        if st.button("Load URL"):
            try:
                r = __import__('requests').get(url); r.raise_for_status()
                parsed = urlparse(url); ext = parsed.path.split('.')[-1].lower()
                if ext in ('csv','txt'):
                    df = pd.read_csv(io.StringIO(r.text))
                elif ext in ('json'):
                    df = pd.json_normalize(r.json())
                elif ext in ('xlsx','xls'):
                    df = pd.read_excel(io.BytesIO(r.content))
                else:
                    df = pd.json_normalize(r.json())
                st.session_state.df = df
                st.success("Loaded data from URL.")
            except Exception as e:
                st.error(f"URL load failed: {e}")

    df = st.session_state.df
    if df is not None:
        st.subheader("Preview")
        st.dataframe(df.head(200))
        if st.button("Show dtypes"):
            st.write(df.dtypes)
        if st.button("Save snapshot"):
            p = save_snapshot(df)
            st.success(f"Saved snapshot: {p}")

# ---------------------------
# Cleaning & Prep
# ---------------------------
elif page == "Cleaning & Prep":
    st.header("🧹 Cleaning & Preparation")
    df = st.session_state.df
    if df is None:
        st.info("Load a dataset first.")
    else:
        left, right = st.columns([2,1])
        with left:
            if st.button("Drop duplicates"):
                df = df.drop_duplicates().reset_index(drop=True); st.session_state.df = df; st.success("Dropped duplicates")
            if st.button("Drop all-empty rows"):
                df = df.dropna(how="all"); st.session_state.df = df; st.success("Dropped empty rows")
            fill = st.selectbox("Fill missing values", ["--","mean (numeric)","median (numeric)","mode (cat)","0","ffill"])
            if st.button("Apply fill"):
                if fill != "--":
                    for c in df.columns:
                        if df[c].isnull().any():
                            if fill == "mean (numeric)" and pd.api.types.is_numeric_dtype(df[c]): df[c].fillna(df[c].mean(), inplace=True)
                            elif fill == "median (numeric)" and pd.api.types.is_numeric_dtype(df[c]): df[c].fillna(df[c].median(), inplace=True)
                            elif fill == "mode (cat)": df[c].fillna(df[c].mode().iloc[0] if not df[c].mode().empty else "", inplace=True)
                            elif fill == "0": df[c].fillna(0, inplace=True)
                            elif fill == "ffill": df[c].fillna(method="ffill", inplace=True)
                    st.session_state.df = df; st.success("Filled missing values")
            if st.button("Save cleaned copy to session"):
                st.session_state.cleaned_df = df.copy(); st.success("Saved cleaned copy")
        with right:
            st.subheader("Column operations")
            col = st.selectbox("Select column", df.columns.tolist())
            if col:
                st.write(df[col].describe(include="all"))
                if st.button("Drop column"):
                    df = df.drop(columns=[col]); st.session_state.df = df; st.success(f"Dropped {col}")

# ---------------------------
# Auto-EDA
# ---------------------------
elif page == "Auto-EDA":
    st.header("📊 Auto-EDA / Profiling")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        if ENABLE_EDA and YDATA_AVAILABLE:
            with st.spinner("Running profile..."):
                profile = ProfileReport(df, minimal=True)
                html = profile.to_html()
                st.components.v1.html(html, height=800, scrolling=True)
                b64 = base64.b64encode(html.encode()).decode()
                st.markdown(f'<a href="data:text/html;base64,{b64}" download="eda_report.html">📥 Download EDA (HTML)</a>', unsafe_allow_html=True)
        else:
            st.write(df.describe(include="all"))
            st.write("Missing values:")
            st.write(df.isnull().sum())

# ---------------------------
# Visualizations
# ---------------------------
elif page == "Visualizations":
    st.header("📈 Visualizations")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        numeric = df.select_dtypes(include=[np.number]).columns.tolist()
        categorical = df.select_dtypes(include=["object","category"]).columns.tolist()
        st.subheader("Single column charts")
        col = st.selectbox("Numeric column", numeric) if numeric else None
        if col:
            fig = px.line(df.reset_index(), x=df.index, y=col, title=f"{col} Trend", template="plotly_dark")
            st.plotly_chart(fig, use_container_width=True)
            st.plotly_chart(px.histogram(df, x=col, nbins=30, title=f"{col} distribution"), use_container_width=True)
        st.subheader("Grouped charts")
        if categorical and col:
            group = st.selectbox("Group by", ["--"] + categorical)
            if group and group != "--":
                xaxis = "date" if "date" in df.columns else df.index
                fig = px.line(df, x=xaxis, y=col, color=group, title=f"{col} by {group}")
                st.plotly_chart(fig, use_container_width=True)
        if numeric:
            corr = df[numeric].corr()
            st.plotly_chart(px.imshow(corr, text_auto=True, title="Correlation Heatmap"), use_container_width=True)

# ---------------------------
# Drill-Down (interactive)
# ---------------------------
elif page == "Drill-Down":
    st.header("📊 Drill-Down Dashboard")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        # Build base bar chart by region or top categorical
        cat_cols = df.select_dtypes(include=["object","category"]).columns.tolist()
        if not cat_cols:
            st.warning("No categorical columns for drill-down.")
        else:
            base = st.selectbox("Base dimension", cat_cols)
            metric = st.selectbox("Metric (numeric)", df.select_dtypes(include=[np.number]).columns.tolist())
            fig = px.bar(df.groupby(base)[metric].sum().reset_index(), x=base, y=metric, title=f"{metric} by {base}")
            if PLOTLY_EVENTS_AVAILABLE:
                selected = plotly_events(fig, click_event=True)
                st.plotly_chart(fig, use_container_width=True)
                if selected:
                    val = selected[0].get("x")
                    st.write("Drill-down for:", val)
                    filt = df[df[base] == val]
                    st.plotly_chart(px.line(filt, x="date" if "date" in filt.columns else filt.index, y=metric, title=f"{metric} over time for {val}"), use_container_width=True)
            else:
                st.plotly_chart(fig, use_container_width=True)
                st.info("Install streamlit-plotly-events for click-driven drilldown interaction.")

# ---------------------------
# Outliers
# ---------------------------
elif page == "Outliers":
    st.header("🔍 Outlier Detection")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        numeric = df.select_dtypes(include=[np.number]).columns.tolist()
        if not numeric:
            st.warning("No numeric columns.")
        else:
            col = st.selectbox("Column", numeric)
            method = st.selectbox("Method", ["IQR","Z-score","IsolationForest"])
            if st.button("Detect"):
                s = df[col].dropna()
                try:
                    if method == "IQR":
                        Q1,Q3 = s.quantile(0.25), s.quantile(0.75); IQR = Q3-Q1
                        mask = (s < (Q1 - 1.5*IQR)) | (s > (Q3 + 1.5*IQR))
                        out = df.loc[mask.index[mask]]
                    elif method == "Z-score":
                        z = (s - s.mean())/s.std(); mask = z.abs() > 3; out = df.loc[mask.index[mask]]
                    else:
                        iso = IsolationForest(contamination=0.05, random_state=42)
                        preds = iso.fit_predict(s.values.reshape(-1,1))
                        out = df.iloc[np.where(preds==-1)[0]]
                    st.dataframe(out.head(200)); st.success(f"Found {len(out)} outliers")
                except Exception as e:
                    st.error(f"Error: {e}")

# ---------------------------
# Forecasting
# ---------------------------
elif page == "Forecasting":
    st.header("🔮 Forecasting")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        numeric = df.select_dtypes(include=[np.number]).columns.tolist()
        if not numeric:
            st.warning("No numeric columns.")
        else:
            fc_col = st.selectbox("Numeric column", numeric)
            date_candidates = [c for c in df.columns if "date" in c.lower() or pd.api.types.is_datetime64_any_dtype(df[c])]
            date_col = st.selectbox("Date column (optional)", ["--"] + date_candidates)
            periods = st.number_input("Periods to forecast", min_value=1, max_value=365, value=30)
            if st.button("Run Forecast"):
                try:
                    if date_col != "--":
                        data = df[[date_col, fc_col]].dropna().rename(columns={date_col:"ds", fc_col:"y"})
                        data["ds"] = pd.to_datetime(data["ds"])
                    else:
                        data = df[[fc_col]].dropna().reset_index().rename(columns={"index":"ds", fc_col:"y"})
                    if PROPHET_AVAILABLE:
                        m = Prophet(); m.fit(data)
                        future = m.make_future_dataframe(periods=periods)
                        fc = m.predict(future)
                        fig = go.Figure()
                        fig.add_trace(go.Scatter(x=data["ds"], y=data["y"], name="history"))
                        fig.add_trace(go.Scatter(x=fc["ds"], y=fc["yhat"], name="forecast"))
                        st.plotly_chart(fig, use_container_width=True)
                    else:
                        st.warning("Prophet not installed; install prophet for forecasting.")
                except Exception as e:
                    st.error(f"Forecast error: {e}\n{traceback.format_exc()}")

# ---------------------------
# AutoML
# ---------------------------
elif page == "AutoML":
    st.header("🤖 AutoML")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        target = st.selectbox("Target column", ["--"] + df.columns.tolist())
        if target and target != "--":
            X = df.drop(columns=[target]).select_dtypes(include=[np.number])
            y = df[target]
            if X.shape[1] == 0:
                st.warning("No numeric features.")
            else:
                test_size = st.slider("Test size", 0.1, 0.5, 0.2)
                if st.button("Train RF"):
                    try:
                        stratify = y if not pd.api.types.is_numeric_dtype(y) and len(y.unique())>1 else None
                        X_train,X_test,y_train,y_test = train_test_split(X,y,test_size=test_size,random_state=42,stratify=stratify)
                        if pd.api.types.is_numeric_dtype(y):
                            model = RandomForestRegressor(n_estimators=150, random_state=42)
                            model.fit(X_train,y_train); preds = model.predict(X_test)
                            st.success(f"RMSE: {mean_squared_error(y_test,preds,squared=False):.3f}, R2: {r2_score(y_test,preds):.3f}")
                            st.bar_chart(pd.Series(model.feature_importances_, index=X.columns).nlargest(15))
                            st.session_state.models["automl"] = ("regression", model, X.columns.tolist())
                        else:
                            model = RandomForestClassifier(n_estimators=150, random_state=42)
                            model.fit(X_train,y_train); preds = model.predict(X_test)
                            st.success(f"Accuracy: {accuracy_score(y_test,preds):.3f}")
                            st.write(classification_report(y_test,preds))
                            st.session_state.models["automl"] = ("classification", model, X.columns.tolist())
                    except Exception as e:
                        st.error(f"AutoML error: {e}\n{traceback.format_exc()}")
        if "automl" in st.session_state.models:
            if st.button("Download model"):
                mtype,m,fcols = st.session_state.models["automl"]
                buf = io.BytesIO(); joblib.dump(m, buf); buf.seek(0)
                download_bytes(buf.read(), f"automl_{mtype}.joblib", mime="application/octet-stream")

# ---------------------------
# AI Agent (Insights, Chat, SQL, Charts)
# ---------------------------
elif page == "AI Agent":
    st.header("🧠 AI Agent")
    df = st.session_state.df
    if not ENABLE_AI:
        st.info("Enable AI in sidebar.")
    elif groq_client is None:
        st.error("Groq not configured. Install groq and set GROQ_API_KEY.")
    else:
        st.subheader("AI Insights")
        cols = st.multiselect("Columns to include (empty=all)", df.columns.tolist() if df is not None else [])
        instr = st.text_area("Instructions", value="You are an expert data analyst. Provide key trends, anomalies, and 3 recommendations.")
        if st.button("Generate Insights"):
            usecols = cols if cols else (df.columns.tolist() if df is not None else [])
            summary = df[usecols].describe().to_string() if df is not None else "No data"
            missing = df[usecols].isnull().sum().to_string() if df is not None else ""
            prompt = f"{instr}\n\nSummary:\n{summary}\n\nMissing:\n{missing}\n\nAnswer in bullets."
            with st.spinner("Calling Groq..."):
                ans = groq_chat(prompt)
                st.markdown("#### AI Insights")
                st.write(ans)

        st.subheader("AI Chat")
        chat = st.text_area("Ask about the dataset", height=150)
        if st.button("Send Chat"):
            ctx = f"Preview:\n{df.head(5).to_string()}\nColumns: {', '.join(df.columns.tolist())}\n\n" if df is not None else ""
            prompt = ctx + chat
            with st.spinner("Calling Groq..."):
                res = groq_chat(prompt)
                st.markdown("#### Response")
                st.write(res)

        st.subheader("NL → SQL")
        nl = st.text_input("Describe query (NL)")
        if st.button("Generate SQL"):
            prompt = f"You have table data_table columns: {', '.join(df.columns.tolist()) if df is not None else ''}. Write SQL: {nl}. Return only SQL."
            out = groq_chat(prompt)
            st.code(out)
            if st.button("Run SQL (requires SQLite load)"):
                if st.session_state.sqlite_conn is None:
                    st.warning("Load dataset to SQLite first on SQL Engine page.")
                else:
                    try:
                        res = pd.read_sql_query(out, st.session_state.sqlite_conn)
                        st.dataframe(res.head(300))
                    except Exception as e:
                        st.error(f"SQL execution error: {e}")

        st.subheader("NL → Chart code")
        chart_desc = st.text_input("Describe chart")
        if st.button("Generate Chart Code"):
            prompt = f"Write python plotly express code that given df creates chart: {chart_desc}. Return just code and set variable 'fig'."
            code = groq_chat(prompt)
            st.code(code)
            if SAFE_EXEC and CONFIRM_EXEC:
                try:
                    local_env = {"df": df.copy(), "px": px, "pd": pd, "np": np}
                    exec(code, {}, local_env)
                    fig = local_env.get("fig")
                    if fig is not None:
                        st.plotly_chart(fig, use_container_width=True)
                    else:
                        st.warning("No 'fig' variable returned.")
                except Exception as e:
                    st.error(f"Execution error: {e}")
            else:
                st.info("Enable code execution and confirm to run generated code.")

# ---------------------------
# AI Formula Builder
# ---------------------------
elif page == "AI Formula Builder":
    st.header("✏️ AI Formula Builder")
    df = st.session_state.df
    st.markdown("Type formulas like: `profit_margin = (revenue - cost) / revenue`")
    formula = st.text_input("Formula")
    def apply_formula(df_local, formula_str):
        import re
        try:
            left,right = formula_str.split("=",1)
            col = left.strip()
            expr = right.strip()
            tokens = set(re.findall(r"[A-Za-z_]\w*", expr))
            for t in tokens:
                if t in df_local.columns:
                    expr = re.sub(rf"\b{t}\b", f"df_local['{t}']", expr)
            df_local[col] = eval(expr)
            return df_local, f"Created column {col}"
        except Exception as e:
            return df_local, f"Error: {e}"
    if st.button("Preview formula"):
        if df is None:
            st.error("Load dataset first.")
        else:
            preview, msg = apply_formula(df.copy(), formula)
            st.write(msg)
            st.dataframe(preview.head())
            if st.button("Apply formula to dataset"):
                st.session_state.df = preview
                st.success("Applied formula.")

# ---------------------------
# AI Data Cleaner
# ---------------------------
elif page == "AI Data Cleaner":
    st.header("🧹 AI Data Cleaner")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        st.write("Sample of data:")
        st.dataframe(df.head(10))
        if st.button("Ask AI to suggest cleaning steps"):
            if not ENABLE_AI or groq_client is None:
                st.error("Enable AI and configure Groq first.")
            else:
                prompt = "You are a data engineer. Here are sample rows:\n" + df.head(10).to_csv(index=False) + "\nList JSON: actions with {col,issue,suggestion,pandas_code}."
                out = groq_chat(prompt)
                st.write("AI suggestions (raw):")
                st.write(out)
                st.info("Review suggestions and choose apply.")

# ---------------------------
# Semantic Search
# ---------------------------
elif page == "Semantic Search":
    st.header("🔎 Semantic Search (vector)")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        if not STRANS_AVAILABLE:
            st.warning("sentence-transformers or faiss not installed. Install them to use this.")
        else:
            if st.button("Build embeddings (FAISS)"):
                texts = df.astype(str).agg(' | '.join, axis=1).tolist()
                model = SentenceTransformer('all-MiniLM-L6-v2')
                embs = model.encode(texts, convert_to_numpy=True)
                dim = embs.shape[1]
                idx = faiss.IndexFlatL2(dim)
                idx.add(embs)
                st.session_state.emb_index = idx
                st.session_state.emb_texts = texts
                st.session_state.faiss_dim = dim
                st.success("Built FAISS index.")
            q = st.text_input("Search query (natural language)")
            if st.button("Search"):
                if st.session_state.emb_index is None:
                    st.error("Build embeddings first.")
                else:
                    model = SentenceTransformer('all-MiniLM-L6-v2')
                    qv = model.encode([q])
                    D,I = st.session_state.emb_index.search(qv, k=10)
                    res = df.iloc[I[0]].copy()
                    st.dataframe(res.head(50))

# ---------------------------
# Notebook Mode (Ace)
# ---------------------------
elif page == "Notebook Mode":
    st.header("📓 Notebook Mode (Ace)")
    if not ACE_AVAILABLE:
        st.warning("Install streamlit-ace to enable Notebook Mode.")
    else:
        code = st_ace(value="# Python code\n# Use variable df\nprint(df.head())", language='python', theme='monokai', height=300)
        if st.button("Run cell"):
            try:
                local_env = {"df": st.session_state.df.copy() if st.session_state.df is not None else pd.DataFrame(), "pd": pd, "np": np, "px": px}
                exec(code, {}, local_env)
                st.success("Executed. Inspect local_env for variables.")
                if "result" in local_env:
                    st.write(local_env["result"])
            except Exception as e:
                st.error(f"Execution error: {e}")

# ---------------------------
# Feature Engineering Studio
# ---------------------------
elif page == "Feature Studio":
    st.header("🧬 Feature Engineering Studio")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        numeric = df.select_dtypes(include=[np.number]).columns.tolist()
        if st.button("Create lag features"):
            lag = st.number_input("Lag steps", min_value=1, max_value=30, value=1)
            for c in numeric:
                df[f"{c}_lag{lag}"] = df[c].shift(lag)
            st.session_state.df = df
            st.success("Lag features created.")
        if st.button("Create rolling mean"):
            window = st.number_input("Window", min_value=2, max_value=100, value=7)
            for c in numeric:
                df[f"{c}_rolling{window}"] = df[c].rolling(window).mean()
            st.session_state.df = df
            st.success("Rolling features created.")

# ---------------------------
# DAX-like Measures (simple)
# ---------------------------
elif page == "DAX Measures":
    st.header("🧮 DAX-like Measures (simple)")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        measure = st.text_input("Define measure (e.g. Total_Revenue = SUM(units_sold * price))")
        if st.button("Compute measure"):
            try:
                left,right = measure.split("=",1)
                name = left.strip()
                expr = right.strip()
                # very naive: allow SUM(expr) only
                if expr.upper().startswith("SUM(") and expr.endswith(")"):
                    inner = expr[4:-1]
                    val = df.eval(inner).sum()
                    st.metric(name, val)
                else:
                    st.error("Only SUM(...) supported in this simple parser.")
            except Exception as e:
                st.error(f"Parse error: {e}")

# ---------------------------
# Time-Series Anomaly Detection
# ---------------------------
elif page == "TS Anomaly":
    st.header("⏱️ Time-Series Anomaly Detection")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        numeric = df.select_dtypes(include=[np.number]).columns.tolist()
        if not numeric:
            st.warning("No numeric columns.")
        else:
            col = st.selectbox("Column", numeric)
            if st.button("Detect anomalies (ruptures)"):
                if not RUPTURES_AVAILABLE:
                    st.error("Install ruptures for TS anomaly detection.")
                else:
                    series = df[col].fillna(method='ffill').values
                    algo = rpt.Pelt(model="rbf").fit(series)
                    result = algo.predict(pen=10)
                    st.write("Change points (indices):", result)
                    fig = px.line(df, y=col)
                    for r in result:
                        fig.add_vline(x=r, line_dash="dash", line_color="red")
                    st.plotly_chart(fig, use_container_width=True)

# ---------------------------
# Live Connectors
# ---------------------------
elif page == "Live Connectors":
    st.header("🌐 Live Connectors")
    source = st.selectbox("Type", ["REST URL (CSV/JSON/XLSX)", "Google Sheets (public)", "Database (Postgres/MySQL)"])
    if source.startswith("REST"):
        url = st.text_input("Enter URL")
        if st.button("Load"):
            try:
                r = __import__('requests').get(url); r.raise_for_status()
                ext = urlparse(url).path.split('.')[-1].lower()
                if ext in ('csv','txt'):
                    df = pd.read_csv(io.StringIO(r.text))
                elif ext in ('json'):
                    df = pd.json_normalize(r.json())
                elif ext in ('xlsx','xls'):
                    df = pd.read_excel(io.BytesIO(r.content))
                else:
                    df = pd.json_normalize(r.json())
                st.session_state.df = df; st.success("Loaded data.")
            except Exception as e:
                st.error(f"Load error: {e}")
    else:
        st.info("Database connectors & Google Sheets integration can be added (requires credentials).")

# ---------------------------
# Designer (save layout JSON)
# ---------------------------
elif page == "Designer":
    st.header("🎛️ Dashboard Designer (save/load simple layouts)")
    comps = ["kpi","line","bar","table","heatmap"]
    selected = st.multiselect("Pick components", comps)
    slots = [None]*6
    for i, s in enumerate(selected[:6]):
        slots[i] = s
    if st.button("Save layout"):
        layout = {"slots":slots, "timestamp": datetime.now().isoformat()}
        os.makedirs("layouts", exist_ok=True)
        fname = f"layouts/layout_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        with open(fname,"w",encoding="utf-8") as f: json.dump(layout,f)
        st.success(f"Saved {fname}")
    st.write("Current slots:", slots)

# ---------------------------
# SQL Engine
# ---------------------------
elif page == "SQL Engine":
    st.header("🗄️ SQL Engine (SQLite in-memory)")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        if st.button("Load into SQLite"):
            try:
                conn = sqlite3.connect(":memory:")
                df.to_sql("data_table", conn, index=False, if_exists="replace")
                st.session_state.sqlite_conn = conn
                st.success("Loaded into SQLite as data_table")
            except Exception as e:
                st.error(f"Load error: {e}")
        if st.session_state.sqlite_conn:
            sql = st.text_area("SQL query (run on data_table)", value="SELECT * FROM data_table LIMIT 100")
            if st.button("Execute SQL"):
                try:
                    res = pd.read_sql_query(sql, st.session_state.sqlite_conn)
                    st.dataframe(res.head(500))
                except Exception as e:
                    st.error(f"SQL error: {e}")

# ---------------------------
# Export Center (PDF/PPTX/CSV)
# ---------------------------
elif page == "Export Center":
    st.header("📤 Export & Reports")
    if st.session_state.cleaned_df is not None:
        st.download_button("Download cleaned CSV", data=st.session_state.cleaned_df.to_csv(index=False).encode(), file_name="cleaned.csv", mime="text/csv")
        st.download_button("Download cleaned Excel", data=df_to_excel_bytes(st.session_state.cleaned_df), file_name="cleaned.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
    elif st.session_state.df is not None:
        st.download_button("Download dataset CSV", data=st.session_state.df.to_csv(index=False).encode(), file_name="dataset.csv", mime="text/csv")
    if REPORTLAB_AVAILABLE and st.session_state.df is not None:
        if st.button("Generate AI PDF Summary"):
            prompt = f"Summarize dataset columns: {', '.join(st.session_state.df.columns.tolist())}. Top 5 insights and 3 recommendations."
            text = groq_chat(prompt) if ENABLE_AI and groq_client else "AI disabled"
            buf = io.BytesIO(); c = canvas.Canvas(buf, pagesize=letter)
            y = 750; c.setFont("Helvetica-Bold", 14); c.drawString(30,y,"AI Generated Summary"); y-=30; c.setFont("Helvetica",10)
            for line in text.splitlines():
                if y < 40:
                    c.showPage(); y = 750
                c.drawString(30,y,line[:110]); y -= 12
            c.save(); buf.seek(0)
            st.download_button("Download AI PDF", data=buf, file_name="ai_summary.pdf", mime="application/pdf")
    if PPTX_AVAILABLE and st.session_state.df is not None:
        if st.button("Generate PPTX Report"):
            prs = Presentation(); slide = prs.slides.add_slide(prs.slide_layouts[5])
            title = slide.shapes.title; title.text = "AI Executive Summary"
            tx = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5)); tf = tx.text_frame; tf.text = "Generated by Enterprise Dashboard v3"
            out = io.BytesIO(); prs.save(out); out.seek(0)
            st.download_button("Download PPTX", data=out, file_name="ai_report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# ---------------------------
# Data Quality Scoring
# ---------------------------
elif page == "Data Quality":
    st.header("📈 Data Quality Scoring")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        def dq_score(df_local):
            total = len(df_local)*df_local.shape[1]
            missing = df_local.isnull().sum().sum()
            completeness = 1 - missing/total
            uniqueness = np.mean([df_local[c].nunique()/len(df_local) for c in df_local.columns])
            score = 100*(0.6*completeness + 0.4*uniqueness)
            return round(score,1), {"completeness":round(completeness,3), "uniqueness":round(uniqueness,3)}
        score, meta = dq_score(df)
        st.metric("Data Quality Score", f"{score}/100")
        st.json(meta)

# ---------------------------
# Auth & RBAC (basic)
# ---------------------------
elif page == "Auth & RBAC":
    st.header("🔐 Auth & RBAC (Demo)")
    st.markdown("This is a demo. For production integrate OAuth / SSO.")
    users = st.session_state.users
    st.write("Users (demo):", list(users.keys()))
    if st.button("Add demo user 'test'"):
        st.session_state.users["test"] = "test"; st.success("Added demo user 'test'")

# ---------------------------
# Autosave & Versioning
# ---------------------------
elif page == "Autosave & Versioning":
    st.header("💾 Autosave & Versioning")
    df = st.session_state.df
    if df is None:
        st.info("Load dataset first.")
    else:
        if st.button("Save snapshot now"):
            p = save_snapshot(df)
            st.success(f"Saved snapshot: {p}")
        snaps = sorted(os.listdir("snapshots")) if os.path.exists("snapshots") else []
        st.write("Available snapshots:", snaps)
        choice = st.selectbox("Restore snapshot", ["--"] + snaps)
        if st.button("Restore"):
            if choice != "--":
                st.session_state.df = pd.read_parquet(os.path.join("snapshots", choice))
                st.success(f"Restored {choice}")

# ---------------------------
# Streaming (Simulated)
# ---------------------------
elif page == "Streaming (Simulated)":
    st.header("📡 Streaming (Simulated)")
    if st.button("Start simulated stream (50 events)"):
        placeholder = st.empty()
        for i in range(50):
            new = {"ts": pd.Timestamp.now(), "value": float(np.random.randn())}
            placeholder.write(new)
            time.sleep(0.2)
        st.success("Simulation complete.")

# ---------------------------
# Deployment helpers (Dockerfile)
# ---------------------------
elif page == "Deployment":
    st.header("🚀 Deployment Helpers")
    if st.button("Write Dockerfile & README"):
        try:
            docker_text = """FROM python:3.10-slim

WORKDIR /app
COPY . /app

RUN pip install --upgrade pip
RUN pip install -r requirements.txt

EXPOSE 8501
CMD ["streamlit", "run", "enterprise_dashboard_v3.py", "--server.port=8501", "--server.address=0.0.0.0"]
"""
            with open("Dockerfile","w",encoding="utf-8") as f: f.write(docker_text)
            with open("README_DEPLOY.md","w",encoding="utf-8") as f: f.write("Run docker build -t enterprise-ai .")
            st.success("Wrote Dockerfile and README_DEPLOY.md")
        except Exception as e:
            st.error(f"Write error: {e}")

# End of pages
st.markdown("---")
st.caption("Enterprise AI Analytics v3 — single-file. ⚠️ Be careful with AI code execution and API keys.")
